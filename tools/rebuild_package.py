#!/usr/bin/env python3
"""
Rebuild the Momentum / GBF case package into a clean, workstream-oriented
structure. Reading material is delivered as PDF, fill-in templates as Word,
data and trackers as Excel, and the board deck as PowerPoint.

Final structure:
  Participants/
    0_Start_Here/            (PDF)
    1_Client_Data_Room/      (PDF only - no datasets)
    2_Workflow_A_Impact_Strategy/          (PDF + Word templates)
    3_Workflow_B_Operations_and_Analytics/ (PDF + Word + Excel datasets + dashboard)
    4_Shared_Toolkit/        (PDF + Excel toolkit + PowerPoint + Word templates)
  Facilitators/
    PDF/  Word/              (staff only)
  Answer_Key/                  (facilitator only - case key)
"""

import shutil
import subprocess
import sys
import unicodedata
from pathlib import Path

from docx import Document
from docx.document import Document as _DocClass
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from fpdf import FPDF

ROOT = Path(__file__).resolve().parents[1]
PART = ROOT / "Participants"
FAC = ROOT / "Facilitators"
KEY = ROOT / "Answer_Key"
TOOLS = ROOT / "tools"
CACHE = TOOLS / ".cache"
VENV_PY = ROOT / ".venv" / "bin" / "python"

sys.path.insert(0, str(TOOLS))


# ---------------------------------------------------------------------------
# Text sanitisation for the core PDF font (Helvetica is Latin-1 only)
# ---------------------------------------------------------------------------

def _sanitize(text: str) -> str:
    if not text:
        return ""
    repl = {
        "\u2014": "-", "\u2013": "-", "\u2018": "'", "\u2019": "'",
        "\u201c": '"', "\u201d": '"', "\u2022": "-", "\u2192": "->",
        "\u2265": ">=", "\u2264": "<=", "\u2026": "...",
    }
    for k, v in repl.items():
        text = text.replace(k, v)
    return unicodedata.normalize("NFKD", text).encode("ascii", "ignore").decode("ascii")


# ---------------------------------------------------------------------------
# PDF renderer
# ---------------------------------------------------------------------------

class PdfExport(FPDF):
    def footer(self):
        self.set_y(-12)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(120, 120, 120)
        self.cell(0, 8, "Momentum Programme  |  Scaling with Integrity  |  "
                  f"Page {self.page_no()}", align="C")

    def add_title(self, text):
        self.ln(1)
        self.set_font("Helvetica", "B", 17)
        self.set_text_color(27, 42, 74)
        self.multi_cell(0, 8, _sanitize(text))
        self.set_draw_color(0, 122, 135)
        self.set_line_width(0.6)
        y = self.get_y() + 1
        self.line(self.l_margin, y, self.w - self.r_margin, y)
        self.ln(4)

    def add_band(self, text):
        self.set_font("Helvetica", "B", 10)
        self.set_text_color(0, 122, 135)
        self.multi_cell(0, 5, _sanitize(text))
        self.ln(1)

    def add_heading(self, text):
        self.ln(1)
        self.set_font("Helvetica", "B", 12)
        self.set_text_color(27, 42, 74)
        self.multi_cell(0, 6, _sanitize(text))
        self.ln(0.5)

    def add_para(self, text, italic=False):
        if not text.strip():
            self.ln(2)
            return
        self.set_font("Helvetica", "I" if italic else "", 10)
        self.set_text_color(45, 45, 45)
        self.multi_cell(0, 5, _sanitize(text))
        self.ln(1)

    def add_list_item(self, text, prefix="-  "):
        self.set_font("Helvetica", "", 10)
        self.set_text_color(45, 45, 45)
        x = self.l_margin + 4
        self.set_x(x)
        self.multi_cell(self.w - self.r_margin - x, 5, _sanitize(prefix + text))
        self.ln(0.5)

    def _break_word(self, word, w):
        """Split a single token that is wider than the column into pieces."""
        pieces, cur = [], ""
        for ch in word:
            if self.get_string_width(cur + ch) <= w - 2 or not cur:
                cur += ch
            else:
                pieces.append(cur)
                cur = ch
        if cur:
            pieces.append(cur)
        return pieces or [word]

    def _wrap(self, text, w):
        text = _sanitize(text)
        words = text.split()
        lines, cur = [], ""
        for word in words:
            if self.get_string_width(word) > w - 2:
                if cur:
                    lines.append(cur)
                    cur = ""
                pieces = self._break_word(word, w)
                lines.extend(pieces[:-1])
                cur = pieces[-1]
                continue
            trial = (cur + " " + word).strip()
            if self.get_string_width(trial) <= w - 2 or not cur:
                cur = trial
            else:
                lines.append(cur)
                cur = word
        if cur:
            lines.append(cur)
        return lines or [""]

    def add_grid(self, rows, header=True, widths=None):
        if not rows:
            return
        cols = max(len(r) for r in rows)
        epw = self.w - self.l_margin - self.r_margin
        if widths and len(widths) == cols:
            tot = sum(widths)
            col_w = [epw * x / tot for x in widths]
        else:
            col_w = [epw / cols] * cols
        line_h = 4.6
        for ri, row in enumerate(rows):
            cells = [str(c) for c in row] + [""] * (cols - len(row))
            is_header = header and ri == 0
            self.set_font("Helvetica", "B" if is_header else "", 8.5)
            wrapped = [self._wrap(cells[ci], col_w[ci]) for ci in range(cols)]
            nlines = max(len(w) for w in wrapped)
            row_h = line_h * nlines + 1.5
            if self.get_y() + row_h > self.page_break_trigger:
                self.add_page()
            x0, y0 = self.l_margin, self.get_y()
            if is_header:
                self.set_fill_color(27, 42, 74)
                self.set_text_color(255, 255, 255)
            elif ri % 2 == 0:
                self.set_fill_color(242, 244, 247)
                self.set_text_color(40, 40, 40)
            else:
                self.set_fill_color(255, 255, 255)
                self.set_text_color(40, 40, 40)
            for ci in range(cols):
                x = x0 + sum(col_w[:ci])
                self.rect(x, y0, col_w[ci], row_h, style="FD")
                ty = y0 + 0.8
                for ln in wrapped[ci]:
                    self.set_xy(x + 1, ty)
                    self.cell(col_w[ci] - 2, line_h, ln)
                    ty += line_h
            self.set_xy(x0, y0 + row_h)
        self.set_text_color(45, 45, 45)
        self.ln(2)


def _iter_blocks(doc):
    body = doc.element.body
    for child in body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, doc)
        elif isinstance(child, CT_Tbl):
            yield Table(child, doc)


def _para_kind(para):
    style = para.style.name if para.style else ""
    if "Heading 1" in style:
        return "heading"
    if "Heading" in style:
        return "subheading"
    if "List Bullet" in style:
        return "bullet"
    if "List Number" in style:
        return "number"
    runs = [r for r in para.runs if r.text.strip()]
    if runs:
        r0 = runs[0]
        size = r0.font.size.pt if r0.font.size else None
        if r0.bold and size and size >= 20:
            return "title"
        if r0.bold and size and size >= 13:
            return "heading"
        if r0.font.color is not None and r0.font.color.type is not None:
            try:
                if r0.font.color.rgb is not None and str(r0.font.color.rgb) == "007A87":
                    return "band"
            except Exception:
                pass
    return "body"


def _row_is_header(row):
    seen = False
    for cell in row.cells:
        text = cell.text.strip()
        if not text:
            continue
        seen = True
        bold = any(run.bold for p in cell.paragraphs for run in p.runs)
        if not bold:
            return False
    return seen


def docx_to_pdf(docx_path: Path, pdf_path: Path):
    doc = Document(str(docx_path))
    pdf = PdfExport()
    pdf.set_auto_page_break(auto=True, margin=16)
    pdf.set_margins(18, 16, 18)
    pdf.add_page()
    number_counter = 0
    prev_kind = None
    for block in _iter_blocks(doc):
        if isinstance(block, Paragraph):
            text = block.text.strip()
            if not text:
                continue
            kind = _para_kind(block)
            if kind != "number":
                number_counter = 0
            if kind == "title":
                pdf.add_title(text)
            elif kind == "band":
                pdf.add_band(text)
            elif kind == "heading":
                pdf.add_heading(text)
            elif kind == "subheading":
                pdf.add_heading(text)
            elif kind == "bullet":
                pdf.add_list_item(text)
            elif kind == "number":
                number_counter += 1
                pdf.add_list_item(text, prefix=f"{number_counter}.  ")
            else:
                italic = any(r.italic for r in block.runs if r.text.strip())
                pdf.add_para(text, italic=italic)
            prev_kind = kind
        else:  # Table
            rows = [[c.text.strip() for c in r.cells] for r in block.rows]
            rows = [r for r in rows if any(cell for cell in r)]
            if not rows:
                continue
            header = _row_is_header(block.rows[0])
            widths = []
            for c in block.rows[0].cells:
                widths.append(c.width.inches if c.width else None)
            if any(w is None or w <= 0 for w in widths):
                widths = None
            pdf.add_grid(rows, header=header, widths=widths)
    pdf_path.parent.mkdir(parents=True, exist_ok=True)
    pdf.output(str(pdf_path))


# ---------------------------------------------------------------------------
# Excel: datasets workbook
# ---------------------------------------------------------------------------

def csv_to_excel(raw_dir: Path, out_xlsx: Path):
    import csv
    wb = Workbook()
    wb.remove(wb.active)
    header_fill = PatternFill("solid", fgColor="007A87")
    header_font = Font(name="Calibri", bold=True, color="FFFFFF")

    ws = wb.create_sheet("DATA_DICTIONARY")
    ws["A1"] = "GBF datasets - field notes and known data-quality issues"
    ws["A1"].font = Font(bold=True, size=13, color="1B2A4A")
    ws["A2"] = "See the Workflow B brief (DB-1) and WSB_Data_Dictionary.pdf. These issues are intentional; find and handle them."
    ws["A2"].font = Font(italic=True, color="666666")
    ws["A4"], ws["B4"] = "Field / area", "What to watch for"
    for c in ("A4", "B4"):
        ws[c].font = header_font
        ws[c].fill = header_fill
    notes = [
        ("beneficiary_id", "Duplicate IDs exist by design; de-duplicate carefully"),
        ("placed_90d", "Coded as Y / Yes / 1 / N / No / 0 / blank - normalise before counting"),
        ("province", "'HCMC' and 'Ho Chi Minh City' both appear - standardise"),
        ("date fields", "Mixed formats (dd/mm/yyyy and yyyy-mm-dd)"),
        ("age", "At least one impossible value - investigate, do not delete silently"),
        ("wage_vnd_monthly", "At least one obvious typo/outlier"),
        ("hub = DN (Dong Nai)", "Underperforms; Q3 costs appear double-entered in hub_costs_2025"),
        ("volunteer_hours", "Hours logged against an inactive mentor - exclude for DB-4"),
        ("employment_type", "Formal vs Gig - the placement rate depends on which you count"),
        ("sample size", "This export (~520) is smaller than the 2,847 in D-01; state this limitation"),
    ]
    for i, (a, b) in enumerate(notes, 5):
        ws.cell(row=i, column=1, value=a)
        ws.cell(row=i, column=2, value=b)
    ws.column_dimensions["A"].width = 26
    ws.column_dimensions["B"].width = 70

    for csv_file in sorted(raw_dir.glob("*.csv")):
        sheet = wb.create_sheet(csv_file.stem[:31])
        with open(csv_file, newline="", encoding="utf-8") as f:
            for ri, row in enumerate(csv.reader(f), 1):
                for ci, val in enumerate(row, 1):
                    cell = sheet.cell(row=ri, column=ci, value=val)
                    if ri == 1:
                        cell.font = header_font
                        cell.fill = header_fill
        for col in sheet.columns:
            sheet.column_dimensions[col[0].column_letter].width = 16
        sheet.freeze_panes = "A2"
    out_xlsx.parent.mkdir(parents=True, exist_ok=True)
    wb.save(out_xlsx)
    print(f"  Excel: {out_xlsx.relative_to(ROOT)}")


# ---------------------------------------------------------------------------
# PowerPoint board deck template
# ---------------------------------------------------------------------------

def build_ppt(out_path: Path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    NAVY = RGBColor(0x1B, 0x2A, 0x4A)
    TEAL = RGBColor(0x00, 0x7A, 0x87)
    GREY = RGBColor(0x66, 0x66, 0x66)

    def title_slide():
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.6))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "GBF Strategy 2026-2028"
        p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = NAVY
        p.font.name = "Lexend"
        p2 = tf.add_paragraph()
        p2.text = "[Your one-sentence recommendation goes here]"
        p2.font.size = Pt(20); p2.font.color.rgb = TEAL
        p2.font.name = "Nunito"
        p3 = tf.add_paragraph()
        p3.text = "Momentum Programme  |  Generation Bridge Foundation  |  13 August 2026"
        p3.font.size = Pt(13); p3.font.color.rgb = GREY
        p3.font.name = "Nunito"

    def content_slide(headline, evidence, insight, ask, notes=""):
        s = prs.slides.add_slide(blank)
        hb = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.3))
        htf = hb.text_frame; htf.word_wrap = True
        hp = htf.paragraphs[0]
        hp.text = headline
        hp.font.size = Pt(26); hp.font.bold = True; hp.font.color.rgb = NAVY
        hp.font.name = "Lexend"
        bb = s.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(12.1), Inches(4.8))
        btf = bb.text_frame; btf.word_wrap = True
        for label, val in [("Evidence", evidence), ("Insight", insight),
                           ("What we recommend / ask", ask)]:
            lp = btf.add_paragraph()
            lr = lp.add_run(); lr.text = label; lr.font.bold = True
            lr.font.size = Pt(15); lr.font.color.rgb = TEAL
            lr.font.name = "Lexend"
            vp = btf.add_paragraph()
            vr = vp.add_run(); vr.text = val
            vr.font.size = Pt(15); vr.font.color.rgb = RGBColor(0x30, 0x30, 0x30)
            vr.font.name = "Nunito"
            btf.add_paragraph()
        if notes:
            s.notes_slide.notes_text_frame.text = notes

    title_slide()
    content_slide(
        "GBF should [action] - because [reason]",
        "[Lead stat + source, e.g. waitlist 2,400; placement 63% formal]",
        "[The single most important 'so what']",
        "[The specific decision you want from the Board]",
        "Open with the answer. You have 60 seconds of Board attention.")
    content_slide(
        "Demand has outpaced funding",
        "Waitlist ~2,400 (D-01); income grew 3% while need grew far faster (D-12)",
        "There is moral urgency to grow - but scaling a broken model repeats Dong Nai",
        "[Frame the central trade-off]")
    content_slide(
        "Young people drop out when transport costs bite",
        "BN-02 and ST-05; survey barrier data; Dong Nai internship drop-out 38% (D-02)",
        "A modest transport stipend may lift completion more than any new hub",
        "[Recommend a transport pilot? Quantify it]")
    content_slide(
        "The headline placement rate flatters us",
        "71% broad vs ~63% formal (ST-02); definition footnote in D-01",
        "Honest dual reporting protects the VPBank relationship, not endangers it",
        "[Recommend a reporting standard]")
    content_slide(
        "We weighed [X] against [Y] and [Z]",
        "Prioritisation matrix; funding scenarios (DB-5)",
        "[The decisive trade-off between the options]",
        "[Name your recommended option]")
    content_slide(
        "A 90-day plan to start well",
        "EP-02 and ST-05 on employer quality; D-11 on staff workload",
        "Early, visible wins build Board and donor confidence",
        "0-30 days: [ ...]   31-60: [ ...]   61-90: [ ...]")
    content_slide(
        "The decision we need today",
        "[Recap the one-line recommendation]",
        "[Why now]",
        "1) Approve direction  2) Authorise pilot budget  3) Mandate honest M&E")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"  PowerPoint: {out_path.relative_to(ROOT)}")


# ---------------------------------------------------------------------------
# Facilitator materials
# ---------------------------------------------------------------------------

def build_facilitator_docs():
    from brand import (set_doc_defaults, add_cover, add_heading, add_body,
                       add_bullet, add_number, add_table, add_quote)

    fac_word = FAC / "Word"
    fac_word.mkdir(parents=True, exist_ok=True)

    def newdoc(title, subtitle):
        d = Document()
        set_doc_defaults(d)
        add_cover(d, title, subtitle, classification="FACILITATOR ONLY - do not distribute",
                  doc_ref="Facilitator pack")
        return d

    # Facilitator Guide
    d = newdoc("Facilitator Guide", "How to run the six-week Momentum engagement")
    add_heading(d, "Before you start", 1)
    add_body(d, "Distribute only the Participants/ folder. Never share Facilitators/ or "
             "Answer_Key/. Confirm each team has five or six members and assign one Buddy per team. "
             "Operational datasets live only in Workflow B's folder - do not hand A the Excel.")
    add_heading(d, "What this case is testing (staff only)", 1)
    add_body(d, "Do not state these lenses to participants.")
    add_bullet(d, "Efficiency versus ethics: growth/targets/placement metrics against dignity, "
               "quality, and who gets left behind. Teams should discover and sit with the tension.")
    add_bullet(d, "Empathy versus outside-in: whether recommendations are grounded in beneficiary "
               "voices or designed for people whose context the team has not taken seriously.")
    add_bullet(d, "Active querying: A must work out what to ask; B must translate. Penalise "
               "teams that reverse-engineer from D-01 alone or share the raw file.")
    add_heading(d, "Kickoff (Day 1) - what to say", 1)
    add_body(d, "Set the tone in the first ten minutes:")
    add_bullet(d, "You are consultants, not students. The Board is paying for judgement.")
    add_bullet(d, "There is no single right answer. We assess your reasoning and evidence.")
    add_bullet(d, "Everything you need is in the pack. Do not go looking outside it.")
    add_bullet(d, "Two workflows hold incomplete information. A frames questions; B owns the data "
               "and translates findings. You must talk to each other.")
    add_heading(d, "Week-by-week facilitation", 1)
    add_table(d, ["Week", "Watch for", "Nudge if..."], [
        ("1", "Vague problem statements; skipping credibility check; A asking for 'the data file'",
         "they treat every source as equally reliable, or A tries to bypass B"),
        ("2", "No Analysis Requests; B answering unasked questions; empathy maps with no citations",
         "Request Log is empty by Wednesday"),
        ("3", "Options that are not genuinely distinct; numbers without Request IDs",
         "all three options are really the same idea"),
        ("4", "Recommendations that ignore Findings or validation",
         "nothing changed after FM-01 or the transcripts"),
        ("5", "Data-dump slides; burying bad news; no named cost of the preferred path",
         "Dong Nai or the formal rate is missing; trade-off reflection is empty"),
        ("6", "Shallow reflection; blame in peer review",
         "the debrief avoids the costs of their recommendation"),
    ], widths=[0.6, 3.0, 3.0])
    add_heading(d, "Escalation triggers", 1)
    add_bullet(d, "Team conflict unresolved for more than 48 hours.")
    add_bullet(d, "Any safeguarding or ethical red flag in how a team behaves.")
    add_bullet(d, "A Buddy giving answers, the data file, or named 'efficiency vs ethics' framing.")
    d.save(fac_word / "01_Facilitator_Guide.docx")

    # Buddy Guide
    d = newdoc("Buddy Guide", "Progressive hints, not answers")
    add_body(d, "Your job is to ask better questions, not to solve the case. Never name "
             "'efficiency versus ethics' or 'elitist outside-in' to participants. Use the "
             "three-level hint ladder. Only move to the next level if the team is genuinely stuck.")
    add_heading(d, "The hint ladder", 1)
    for topic, l1, l2, l3 in [
        ("Problem framing",
         "What exactly is the Board asking you to decide?",
         "Is your problem statement about programmes, geography, funding, or all three?",
         "Ask which voices in the data room pull hardest in opposite directions."),
        ("Asking for evidence",
         "What would you need to know to choose between your options?",
         "Have you filed an Analysis Request, or are you guessing from the annual report?",
         "Ask A to turn one hypothesis into R-0x before inventing a number."),
        ("Placement rate",
         "What does 'placed' actually mean here?",
         "Whose definition are you using - D-01's or ST-02's?",
         "Have B compute both the broad and the formal rate and return them via a Findings Memo."),
        ("Recommendations",
         "What is the evidence for that claim?",
         "Which transcript or Findings Memo backs this up?",
         "Ask them to map every recommendation to a source in the validation protocol, including R-/FM- IDs."),
        ("Costs of the choice",
         "Who gains under your recommendation?",
         "Who waits longer, or is deprioritised, if you choose this path?",
         "Point them to the trade-off reflection without naming the dilemma for them."),
    ]:
        add_heading(d, topic, 2)
        add_number(d, "Level 1: " + l1)
        add_number(d, "Level 2: " + l2)
        add_number(d, "Level 3: " + l3)
    add_heading(d, "Check-in schedule", 1)
    add_body(d, "Fridays of Weeks 1-5. Fifteen minutes. Review Request Log health, Findings Memo "
             "quality, and whether claims are cited. Do not review answers for correctness. "
             "Do not give A the dataset.")
    add_heading(d, "Common mistakes to expect", 1)
    add_bullet(d, "Accepting the 71% placement rate at face value.")
    add_bullet(d, "A reverse-engineering strategy from D-01 without querying B.")
    add_bullet(d, "B sharing the raw workbook or answering questions A never asked.")
    add_bullet(d, "Recommending national scale without fixing Dong Nai.")
    add_bullet(d, "Treating the board email (D-08) as fact.")
    add_bullet(d, "Claiming attendance causes placement.")
    add_bullet(d, "Empty trade-off reflection; no named cost of the preferred path.")
    d.save(fac_word / "02_Buddy_Guide.docx")

    # Assessment Rubric
    d = newdoc("Assessment Rubric", "100 points - descriptors and point bands per level")
    add_body(d, "Score each criterion against the descriptors. Point bands below make Excellent / "
             "Good / Satisfactory / Poor concrete. These bands are expected scales; Core Team may "
             "announce refinements without changing the criterion weights without notice.")
    add_heading(d, "How to convert a level to points", 2)
    add_body(d, "For a criterion worth P points: Excellent = 85-100% of P; Good = 70-84%; "
             "Satisfactory = 50-69%; Poor = 0-49%. Round to the nearest whole point.")
    add_table(d, ["Criterion max (P)", "Excellent", "Good", "Satisfactory", "Poor"], [
        ("15", "13-15", "11-12", "8-10", "0-7"),
        ("20", "17-20", "14-16", "10-13", "0-9"),
        ("10", "9-10", "7-8", "5-6", "0-4"),
        ("5", "5", "4", "3", "0-2"),
    ], widths=[1.4, 1.2, 1.2, 1.4, 1.2])
    for crit, pts, exc, good, sat, poor in [
        ("Problem understanding", 15,
         "Precise, MECE issue tree; hypotheses imply clear asks for B",
         "Solid framing; minor gaps",
         "Generic framing; some overlap",
         "Misreads the problem"),
        ("Research quality", 10,
         "Sources triangulated; citations include section/speaker; contradictions resolved",
         "Most sources used well",
         "Uses sources uncritically or cites code only",
         "Ignores or misuses sources"),
        ("Design thinking", 15,
         "Evidence-based empathy and journey maps grounded in specific voices",
         "Competent maps and ideation",
         "Mechanical, thin on evidence",
         "Superficial or missing"),
        ("Data analysis and translation", 20,
         "Clean data; honest definitions; Findings Memos in plain language; limitations stated",
         "Mostly sound; small errors",
         "Basic analysis; weak translation to A",
         "Errors, raw-file sharing, or misleading charts"),
        ("Recommendations", 20,
         "Distinct options; Request IDs cited; costs of preferred path named",
         "Reasonable, mostly supported",
         "Vague or weakly supported",
         "Unsupported or infeasible"),
        ("Presentation", 10,
         "Answer-first storyline; every slide earns its place",
         "Clear and mostly disciplined",
         "Some data-dumping",
         "Disorganised or decorative"),
        ("Project management", 5,
         "Request Log, RACI, file naming, version control all maintained",
         "Mostly maintained",
         "Patchy",
         "Absent"),
        ("Teamwork and cross-workflow exchange", 5,
         "Genuine request/findings exchange; online handoff confirmed in good faith",
         "Good collaboration",
         "Uneven contribution or one-way dump",
         "Free-riding or A/B bypass"),
    ]:
        add_heading(d, f"{crit} ({pts} points)", 2)
        add_table(d, ["Excellent", "Good", "Satisfactory", "Poor"],
                  [(exc, good, sat, poor)], widths=[1.7, 1.6, 1.6, 1.5])
    add_heading(d, "Sample answers", 2)
    add_body(d, "There is no single model answer. Use Answer_Key/ for expected findings, "
             "solution paths, and deliverable standards. Reward reasoning that names costs, "
             "not a particular destination.")
    d.save(fac_word / "03_Assessment_Rubric.docx")

    # Solution Paths
    d = newdoc("Solution Paths", "CONFIDENTIAL - there is no single correct answer")
    add_body(d, "Reward reasoning, not a particular destination. Strong teams name what their "
             "preferred path costs. Weak answers pick one slogan and ignore the evidence against it.")
    for name, thesis, strengths, risks in [
        ("Path A - Depth before scale",
         "Pause expansion, fix Dong Nai and employer quality, negotiate a smaller quality-focused VPBank grant.",
         "Honest about D-02; addresses EP-02 and ST-05; sustainable.",
         "May leave money on the table; must convince VPBank (see DN-01)."),
        ("Path B - Conditional scale (most common strong answer)",
         "Accept VPBank but gate expansion behind quality milestones; fix Dong Nai first, then Long An.",
         "Balances DN-01 and BD-01; keeps the grant; stages risk.",
         "Execution-heavy; needs credible M&E to hold the line."),
        ("Path C - Employer-led pivot",
         "Certify and audit employers; place only with vetted partners; make quality the brand.",
         "Directly answers EP-02, ST-05, BN-02; differentiates from YouthWorks.",
         "Slower growth; needs employer buy-in."),
        ("Path D - Digital hybrid",
         "Grow Pathway Digital for reach, keep Skills Forward for depth.",
         "Lower unit cost (D-12); serves waitlist faster.",
         "Digital placement is weaker; equity risk (BN-01 has no laptop); only valid with caveats."),
    ]:
        add_heading(d, name, 2)
        add_body(d, "Thesis: " + thesis)
        add_body(d, "Strengths: " + strengths)
        add_body(d, "Risks to probe: " + risks)
    add_heading(d, "Red lines (mark down if crossed)", 1)
    add_bullet(d, "Recommending 60-day reporting as the primary metric to look better.")
    add_bullet(d, "Replacing caseworkers with volunteers (contradicts VR-01).")
    add_bullet(d, "National scale with no fix for Dong Nai.")
    add_bullet(d, "Cherry-picking easier-to-place youth without acknowledging the cost.")
    add_bullet(d, "Workflow A citing 'the dataset' without Request / Findings IDs.")
    add_bullet(d, "Workflow B sharing the raw workbook with A.")
    d.save(fac_word / "04_Solution_Paths.docx")

    # Data Task Solutions
    d = newdoc("Data Task Solutions", "CONFIDENTIAL - expected findings and pitfalls")
    add_body(d, "Workflow B should find these. Workflow A should learn material issues only "
             "through Findings Memos, not by opening the workbook.")
    add_heading(d, "DB-1 Data quality - what B should find", 1)
    add_bullet(d, "Duplicate beneficiary IDs (built in) - de-duplicate before counting.")
    add_bullet(d, "placed_90d coded as Y/Yes/1/N/No/0/blank - normalise.")
    add_bullet(d, "Mixed date formats across tables.")
    add_bullet(d, "A wage outlier (obvious typo) on one record.")
    add_bullet(d, "Dong Nai Q3 costs double-entered in hub_costs_2025.")
    add_bullet(d, "Hours logged against an inactive mentor in volunteer_hours.")
    add_bullet(d, "'HCMC' vs 'Ho Chi Minh City' inconsistency.")
    add_heading(d, "DB-2/DB-3 expected patterns", 1)
    add_bullet(d, "Dong Nai placement clearly below HCMC hubs.")
    add_bullet(d, "Formal-only rate roughly 8 points below the broad rate (aligns with ST-02).")
    add_bullet(d, "Higher attendance associates with placement - but confounded by transport/motivation.")
    add_bullet(d, "Disability placement materially below overall (aligns with ST-03).")
    add_heading(d, "DB-4 / DB-5", 1)
    add_bullet(d, "Mentor hours weakly associate with completion; do not overclaim causation.")
    add_bullet(d, "Aggressive growth struggles to hold 65% honestly without selection pressure - "
               "the efficiency/ethics tension teams should discover.")
    d.save(fac_word / "05_Data_Task_Solutions.docx")

    # Debrief + FAQ
    d = newdoc("Debrief Guide & FAQ", "Run a 90-minute debrief; answer common questions")
    add_heading(d, "Working space and submission (tell teams at kickoff)", 1)
    add_bullet(d, "Shared workspace: Google Drive / Google Workspace folder per team "
               "(or Microsoft 365 if the cohort is assigned that stack).")
    add_bullet(d, "Submission: upload named file to Drive, then submit the link via the weekly "
               "Google Form (Buddy shares URL). Email is backup only.")
    add_bullet(d, "File naming: TeamShortName_Owner_DeliverableCode_Version.ext")
    add_bullet(d, "Package download: GitHub repo -> Code -> Download ZIP.")
    add_heading(d, "Debrief flow (90 minutes)", 1)
    add_number(d, "Team reflection (15 min): what surprised you?")
    add_number(d, "Evidence round (20 min): each team names the finding that changed their mind "
               "(ideally via a Findings Memo, not a slide of pivots).")
    add_number(d, "Trade-off round (20 min): what did your recommendation cost, and for whom? "
               "(Only then may facilitators name efficiency versus ethics if it still has not surfaced.)")
    add_number(d, "Reveal (20 min): show the multiple valid paths; there was no single answer.")
    add_number(d, "Skills forward (15 min): one thing each person will practise next.")
    add_heading(d, "Frequently asked questions", 1)
    add_body(d, "Can Workflow A have the Excel file?", bold=True)
    add_body(d, "No. Asking what to request is part of the exercise. Buddies must not share it.")
    add_body(d, "Can we collect more data or interview real NGO staff?", bold=True)
    add_body(d, "No. GBF and all named parties are fictional for this Impact Case. Everything "
             "required is in the pack. Managing incomplete information is part of the exercise.")
    add_body(d, "How do we confirm the handoff online?", bold=True)
    add_body(d, "Tick the checklist boxes, type full name and date. A Drive comment or sheet "
             "tick is enough - no wet signature.")
    add_body(d, "Which placement rate is correct?", bold=True)
    add_body(d, "Both are 'correct' under their definitions. The learning point is to state the "
             "definition and report honestly (D-01 footnote vs ST-02).")
    add_body(d, "Is there a right recommendation?", bold=True)
    add_body(d, "No. Several paths are defensible (see Solution Paths). Judgement and evidence win.")
    d.save(fac_word / "06_Debrief_Guide_and_FAQ.docx")

    # Presentation Day Plan
    d = newdoc("Presentation Day Plan", "Panel weights, guest brief, and Board format")
    add_heading(d, "1. Score weights on Presentation Day", 1)
    add_body(d, "The 100-point engagement rubric still applies to written work. On Presentation "
             "Day, the live pitch and Q&A are scored with this panel split:")
    add_table(d, ["Panel", "Weight", "Focus"], [
        ("Buddy", "35%", "Process, teamwork, citation discipline, request/findings exchange"),
        ("Core Team (Impact Case / MMT)", "35%", "Substance of recommendation, evidence quality, feasibility"),
        ("Guest Board (invited)", "30%", "Clarity of ask, persuasiveness, composure under challenge"),
    ], widths=[2.2, 1.0, 3.2])
    add_body(d, "Final written report quality is primarily owned by Core Team using the "
             "Assessment Rubric and Answer_Key. Buddy scores inform process marks; guests do "
             "not re-mark the full 100-point written rubric unless Core Team asks them to.")

    add_heading(d, "2. Who judges the final report?", 1)
    add_bullet(d, "Primary: Core Team Impact Case using 03_Assessment_Rubric and Answer_Key.")
    add_bullet(d, "Buddy: process and teamwork evidence (Request Log, handoff ticks, file naming).")
    add_bullet(d, "Guest Board: live pitch only, unless Core Team shares a one-page score sheet.")

    add_heading(d, "3. Briefing guest Board members", 1)
    add_body(d, "Yes - Core Team briefs guests before the session. Guests should not rely on "
             "teams to teach the whole case in the first two minutes.")
    add_number(d, "Send a one-page case brief (Board question, workflows A/B, fictional disclaimer) "
               "48 hours ahead.")
    add_number(d, "15-minute live brief before pitches: what GBF is (fictional), what a strong "
               "ask looks like, red lines (raw data sharing, outside research).")
    add_number(d, "Optional role cards: guests may play ED / Development / Donor / Board hawk - "
               "or sit as a neutral Board. Tell teams which format applies that day.")

    add_heading(d, "4. Format of the final presentation", 1)
    add_body(d, "Default format (recommended):")
    add_bullet(d, "Teams pitch to a simulated GBF Board (10-12 minutes + 8-10 minutes Q&A).")
    add_bullet(d, "Core Team and Buddies sit on the panel; invited guests join as Board members.")
    add_bullet(d, "Optional light role-play: assign 1-2 guests (or Core Team) to named stances "
               "(e.g. growth-oriented Board member vs quality-first ED). Do not require teams "
               "to invent the Board characters themselves.")
    add_bullet(d, "Teams open with their one-sentence recommendation; they do not re-read the "
               "entire data room as 'context'.")
    add_body(d, "Announce the chosen format at Week 4 Buddy check-in so teams can rehearse.")

    add_heading(d, "5. Run-of-show (example)", 1)
    add_table(d, ["Time", "Activity"], [
        ("08:30", "Guest brief (closed)"),
        ("09:00", "Team 1 pitch + Q&A"),
        ("09:25", "Team 2 ..."),
        ("After last team", "Panel scores submitted; short private calibration"),
        ("Closing", "Thank teams; remind Week 6 peer eval and reflection"),
    ], widths=[1.4, 5.0])
    d.save(fac_word / "07_Presentation_Day_Plan.docx")

    # Convert all facilitator docs to PDF
    fac_pdf = FAC / "PDF"
    for docx in sorted(fac_word.glob("*.docx")):
        docx_to_pdf(docx, fac_pdf / docx.with_suffix(".pdf").name)
        print(f"  Facilitator: {docx.stem}")


# ---------------------------------------------------------------------------
# Root quick-start PDFs
# ---------------------------------------------------------------------------

def build_root_pdfs():
    pdf = PdfExport()
    pdf.set_margins(18, 16, 18)
    pdf.add_page()
    pdf.add_title("START HERE")
    pdf.add_band("Momentum Programme  |  Generation Bridge Foundation engagement")
    for line in [
        "Programme runs Monday 13 July to Friday 21 August 2026.",
        "Board presentation: Thursday 13 August 2026, 09:00 ICT.",
        "All final deliverables: Friday 14 August 2026, 17:00 ICT.",
        "",
    ]:
        pdf.add_para(line)
    pdf.add_heading("Read in this order")
    for i, line in enumerate([
        "0_Start_Here/00_Programme_Guideline.pdf  (handbook + weekly deliverables + rules)",
        "0_Start_Here/01_Data_Room_Index.pdf",
        "2_Workflow_A_Impact_Strategy/WSA_Start_Here.pdf",
        "3_Workflow_B_Operations_and_Analytics/WSB_Start_Here.pdf",
        "4_Shared_Toolkit/GBF_Consulting_Toolkit.xlsx (Master Timeline + Request Log)",
    ], 1):
        pdf.add_list_item(line, prefix=f"{i}.  ")
    pdf.add_heading("The rules")
    for line in [
        "Fictional case only - no outside data in deliverables.",
        "Workflow A does not hold the operational datasets; use Analysis Requests.",
        "Workflow B does not share the raw workbook; return Findings Memos.",
        "Cite Ref + section (documents) or CODE (Speaker) for transcripts.",
        "Name files: Team_Owner_DeliverableCode_Version.ext and submit via the weekly Google Form.",
        "Download the pack from GitHub: Code > Download ZIP.",
    ]:
        pdf.add_list_item(line)
    pdf.output(str(ROOT / "START_HERE.pdf"))
    print("  Root: START_HERE.pdf")

    pdf = PdfExport()
    pdf.set_margins(18, 16, 18)
    pdf.add_page()
    pdf.add_title("Package Map")
    pdf.add_heading("Participants/  (distribute this folder)")
    pdf.add_grid([
        ["Folder", "Contents", "Formats"],
        ["0_Start_Here", "Programme Guideline (merged handbook + weekly deadlines), data room index", "PDF"],
        ["1_Client_Data_Room", "GBF documents and interview transcripts (no datasets)", "PDF"],
        ["2_Workflow_A_Impact_Strategy", "Start here, brief, playbook, templates A1-A7", "PDF, Word"],
        ["3_Workflow_B_Operations_and_Analytics", "Start here, brief, dictionary, datasets, dashboard, B templates", "PDF, Word, Excel"],
        ["4_Shared_Toolkit", "Toolkit (Request Log), guides, board deck, shared templates", "Excel, PDF, PowerPoint, Word"],
    ], widths=[2.2, 3.4, 1.6])
    pdf.add_heading("Facilitators/  (staff only - do not share)")
    pdf.add_para("Facilitator guide, buddy guide, assessment rubric (with point bands), "
                 "presentation day plan, solution paths, debrief guide - in PDF and Word.")
    pdf.add_heading("Answer_Key/  (staff only - do not share)")
    pdf.add_para("Case key: data inconsistencies register, expected deliverables and outcomes, "
                 "strategic solution paths, and a searchable Excel register.")
    pdf.output(str(ROOT / "PACKAGE_MAP.pdf"))
    print("  Root: PACKAGE_MAP.pdf")


# ---------------------------------------------------------------------------
# Orchestration
# ---------------------------------------------------------------------------

def clean():
    for target in (PART, FAC, KEY):
        if target.exists():
            shutil.rmtree(target)
    # Remove any stray legacy numbered folders / old readmes
    for p in ROOT.iterdir():
        if p.is_dir() and p.name[:2].isdigit() and p.name not in ("Participants", "Facilitators"):
            shutil.rmtree(p)
    for stray in ("README.md",):
        f = ROOT / stray
        if f.exists():
            f.unlink()


def main():
    print("=== Clean ===")
    clean()

    print("=== Step 1: generate data ===")
    subprocess.run([str(VENV_PY), str(TOOLS / "generate_data.py")], check=True)

    print("=== Step 2: build documents (staging) ===")
    import content
    manifest = content.build_all()

    print("=== Step 3: place documents (PDF read-only / Word templates) ===")
    for item in manifest:
        src = content.STAGING / f"{item['name']}.docx"
        dest_dir = PART / item["dest"]
        if item["fmt"] == "pdf":
            docx_to_pdf(src, dest_dir / f"{item['name']}.pdf")
        else:
            tdir = dest_dir / "Templates"
            tdir.mkdir(parents=True, exist_ok=True)
            shutil.copy(src, tdir / f"{item['name']}.docx")
    print(f"  placed {len(manifest)} documents")

    print("=== Step 4: Excel workbooks ===")
    import generate_excel as ge
    ge.OUT = PART / "4_Shared_Toolkit"
    ge.OUT.mkdir(parents=True, exist_ok=True)
    ge.main()
    csv_to_excel(CACHE / "raw", PART / "3_Workflow_B_Operations_and_Analytics" / "GBF_Datasets.xlsx")

    print("=== Step 5: PowerPoint ===")
    build_ppt(PART / "4_Shared_Toolkit" / "GBF_Board_Deck_Template.pptx")

    print("=== Step 6: Facilitator materials ===")
    build_facilitator_docs()

    print("=== Step 7: Answer Key ===")
    import answer_key as ak
    ak.build_all(docx_to_pdf)

    print("=== Step 8: Root quick-start PDFs ===")
    build_root_pdfs()

    print("\nDone. Structure:")
    for d in sorted(PART.iterdir()):
        print(f"  Participants/{d.name}/")
    print("  Facilitators/PDF, Facilitators/Word")
    print("  Answer_Key/ (PDF, Excel, Word)")


if __name__ == "__main__":
    main()
